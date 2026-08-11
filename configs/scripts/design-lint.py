#!/usr/bin/env python3
"""design-lint — measure rendered design output against the Clade design floors.

The /frontend-design skill declares floors (WCAG 2.2 AA contrast, legibility
minimums, focal-surface caps, presentation word caps).  Most of them are
*rendered-output* rules: a clean grep over source proves nothing about them.
This is the validator that actually opens the artifact and measures.

Three lanes, each usable on its own:

  deck    <file.pptx>        source metrics via python-pptx (optional dep)
  render  <image|dir>        pixel metrics on rendered slides/pages (Pillow)
  html    <file.html|dir>    static analysis of artifact pages (stdlib only)

Exit status is 0 when every check passes, 1 when any FAIL is recorded, and 2
on a usage/environment error.  `--json` emits the full findings for machines.

Honesty contract: every check states what it can and cannot see.  Static HTML
analysis catches definite violations, never proves absence of all of them; the
render lane's contrast probe is a heuristic and is labelled as one.  A check
that cannot run reports SKIP with the reason — it never silently passes.
"""

from __future__ import annotations

import argparse
import json
import re
import sys
from dataclasses import dataclass, field, asdict
from pathlib import Path
from statistics import median
from typing import Any, Iterable

# ─── Defaults (the floors declared in configs/skills/frontend-design) ───

BODY_MIN_PT = 13.0          # hard floor for essential presentation text
BODY_MEDIAN_PT = 18.0       # target median for audience-facing body copy
WORDS_PER_SLIDE = 70        # presentation word cap
FOCAL_COVERAGE_MAX = 0.45   # heavy full-bleed surface share of canvas
CONTRAST_TEXT = 4.5         # WCAG 2.2 AA, normal text
CONTRAST_LARGE = 3.0        # WCAG 2.2 AA, large text / non-text
BODY_MIN_PX = 16.0          # minimum body font-size on the web
INK_DENSITY_MAX = 0.60      # non-background coverage before a slide reads as a wall
TEXT_THINNESS_MIN = 0.30    # boundary/area ratio above which a colour reads as glyphs
TEXT_THINNESS_MAX = 0.98    # at/above this every pixel is an edge — anti-alias fringe
MIN_CANVAS_PX = 400         # smaller images are icons/logos, not design canvases

SEVERITY_ORDER = {"FAIL": 0, "WARN": 1, "SKIP": 2, "PASS": 3}


# ─── Findings ───


@dataclass
class Finding:
    check: str
    severity: str          # FAIL | WARN | SKIP | PASS
    target: str
    detail: str
    measured: Any = None
    threshold: Any = None


@dataclass
class Report:
    lane: str
    findings: list[Finding] = field(default_factory=list)

    def add(self, *a, **kw) -> None:
        self.findings.append(Finding(*a, **kw))

    @property
    def failed(self) -> bool:
        return any(f.severity == "FAIL" for f in self.findings)

    def render(self) -> str:
        rows = sorted(self.findings, key=lambda f: SEVERITY_ORDER.get(f.severity, 9))
        out = []
        for f in rows:
            mark = {"FAIL": "✗", "WARN": "!", "SKIP": "–", "PASS": "✓"}.get(f.severity, "?")
            out.append(f"  {mark} [{f.severity:4}] {f.target}: {f.detail}")
        counts: dict[str, int] = {}
        for f in self.findings:
            counts[f.severity] = counts.get(f.severity, 0) + 1
        summary = " ".join(f"{k}={v}" for k, v in sorted(counts.items()))
        out.append(f"  ── {self.lane}: {summary or 'no checks ran'}")
        return "\n".join(out)


# ─── Colour maths (WCAG 2.1/2.2 relative luminance) ───


def _channel(v: int) -> float:
    c = v / 255.0
    return c / 12.92 if c <= 0.03928 else ((c + 0.055) / 1.055) ** 2.4


def luminance(rgb: tuple[int, int, int]) -> float:
    r, g, b = rgb
    return 0.2126 * _channel(r) + 0.7152 * _channel(g) + 0.0722 * _channel(b)


def contrast_ratio(a: tuple[int, int, int], b: tuple[int, int, int]) -> float:
    la, lb = luminance(a), luminance(b)
    if la < lb:
        la, lb = lb, la
    return (la + 0.05) / (lb + 0.05)


_HEX = re.compile(r"#([0-9a-fA-F]{3}|[0-9a-fA-F]{6})\b")
_RGB = re.compile(r"rgba?\(\s*(\d+)\s*,\s*(\d+)\s*,\s*(\d+)")

_NAMED = {
    "black": (0, 0, 0), "white": (255, 255, 255), "red": (255, 0, 0),
    "green": (0, 128, 0), "blue": (0, 0, 255), "gray": (128, 128, 128),
    "grey": (128, 128, 128), "silver": (192, 192, 192), "navy": (0, 0, 128),
}


_VAR_DEF = re.compile(r"(--[\w-]+)\s*:\s*([^;}]+)")
_VAR_USE = re.compile(r"var\(\s*(--[\w-]+)\s*(?:,\s*([^)]+))?\)")


def custom_properties(css: str) -> dict[str, str]:
    """Collect `--token: value` declarations so var() references can resolve.

    The skill tells designers to build on CSS custom properties, so a checker
    that gives up at `var(--ink)` is blind on exactly the idiom it recommends.
    Later definitions win, matching the cascade for same-specificity rules.
    """
    return {name: value.strip() for name, value in _VAR_DEF.findall(css)}


def resolve_value(value: str, props: dict[str, str], depth: int = 0) -> str:
    """Substitute var() references, honouring fallbacks, with a recursion cap."""
    if depth > 8 or "var(" not in value:
        return value

    def swap(m: re.Match) -> str:
        return props.get(m.group(1), (m.group(2) or "")).strip()

    return resolve_value(_VAR_USE.sub(swap, value), props, depth + 1)


def parse_color(text: str) -> tuple[int, int, int] | None:
    """Best-effort CSS colour → RGB. Returns None when not confidently parsed."""
    text = text.strip().lower()
    if text in _NAMED:
        return _NAMED[text]
    m = _HEX.search(text)
    if m:
        h = m.group(1)
        if len(h) == 3:
            h = "".join(c * 2 for c in h)
        return tuple(int(h[i:i + 2], 16) for i in (0, 2, 4))  # type: ignore[return-value]
    m = _RGB.search(text)
    if m:
        return (int(m.group(1)), int(m.group(2)), int(m.group(3)))
    return None


# ─── Lane: deck (pptx source metrics) ───


def lint_deck(path: Path, report: Report, label: str) -> None:
    try:
        from pptx import Presentation           # type: ignore
        from pptx.util import Inches            # type: ignore
    except ImportError:
        report.add("deck.deps", "SKIP", label,
                   "python-pptx not installed — source metrics unavailable "
                   "(pip install python-pptx); render lane still applies")
        return

    prs = Presentation(str(path))
    canvas_w, canvas_h = prs.slide_width, prs.slide_height

    for idx, slide in enumerate(prs.slides, start=1):
        tag = f"{label} slide {idx}"
        texts: list[str] = []
        body_pt: list[float] = []

        for shape in slide.shapes:
            try:
                left, top = shape.left, shape.top
                width, height = shape.width, shape.height
            except (AttributeError, TypeError):
                left = top = width = height = None
            if None not in (left, top, width, height):
                if left < 0 or top < 0 or left + width > canvas_w or top + height > canvas_h:
                    report.add("deck.bounds", "FAIL", tag,
                               f"shape {shape.shape_id} extends outside the canvas")
            if not getattr(shape, "has_text_frame", False):
                continue
            texts.append(shape.text)
            in_body_zone = top is not None and Inches(2.0) <= top < Inches(6.0)
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    if run.text.strip() and run.font.size and in_body_zone:
                        body_pt.append(run.font.size.pt)

        words = re.findall(r"[A-Za-z0-9$%<+.'’-]+", " ".join(texts))
        if len(words) > WORDS_PER_SLIDE:
            report.add("deck.words", "FAIL", tag,
                       f"{len(words)} words exceeds the presentation cap",
                       measured=len(words), threshold=WORDS_PER_SLIDE)
        else:
            report.add("deck.words", "PASS", tag, f"{len(words)} words",
                       measured=len(words), threshold=WORDS_PER_SLIDE)

        if not body_pt:
            report.add("deck.type", "SKIP", tag,
                       "no explicitly-sized body-zone runs (inherited theme sizes are "
                       "not readable from the file — check the render lane)")
            continue
        lo, mid = min(body_pt), median(body_pt)
        if lo < BODY_MIN_PT:
            report.add("deck.type.floor", "FAIL", tag,
                       f"body-zone text at {lo:.1f}pt is below the legibility floor",
                       measured=round(lo, 1), threshold=BODY_MIN_PT)
        if mid < BODY_MEDIAN_PT:
            report.add("deck.type.median", "WARN", tag,
                       f"body-zone median {mid:.1f}pt is under the presentation target",
                       measured=round(mid, 1), threshold=BODY_MEDIAN_PT)
        if lo >= BODY_MIN_PT and mid >= BODY_MEDIAN_PT:
            report.add("deck.type", "PASS", tag,
                       f"min {lo:.1f}pt / median {mid:.1f}pt")


# ─── Lane: render (pixel metrics on rendered slides/pages) ───


def _pixels(image) -> Iterable:
    """Flat pixel sequence, across the Pillow API rename.

    Pillow 12 introduced `get_flattened_data()` and deprecated `getdata()`
    (removal in Pillow 14).  Deployments run whatever their distro ships, so
    probe rather than assume — picking one name strands the other environment
    with an AttributeError at the exact moment a check should be reporting.
    """
    getter = getattr(image, "get_flattened_data", None) or image.getdata
    return getter()


def _thinness(image, rgb: tuple[int, int, int]) -> float:
    """Fraction of this colour's pixels that lie on its own boundary.

    Glyph strokes are a few pixels wide, so nearly every pixel touches a
    non-matching neighbour and thinness approaches 1.  A filled panel is mostly
    interior and scores near 0.  This is what separates "text that must clear a
    contrast floor" from "a tinted surface that legitimately sits close to the
    page colour".
    """
    from PIL import Image, ImageFilter          # type: ignore

    # The image is already quantised, so exact equality is a safe mask test.
    bits = bytes(255 if px == rgb else 0 for px in _pixels(image))
    mask = Image.frombytes("L", image.size, bits)
    area = sum(1 for v in _pixels(mask) if v)
    if area == 0:
        return 0.0
    interior = mask.filter(ImageFilter.MinFilter(3))
    inner = sum(1 for v in _pixels(interior) if v)
    return 1.0 - (inner / area)


def _local_backdrop(image, rgb: tuple[int, int, int]):
    """The colour this one actually sits on, not the page's dominant colour.

    Text is judged against the surface directly behind it.  A caption on a
    yellow panel over a magenta page must clear its ratio against the yellow;
    comparing it to the page colour scores a slide that is in fact unreadable.
    Returns None when the ring around the colour is not readable.
    """
    from PIL import Image, ImageFilter          # type: ignore

    bits = bytes(255 if px == rgb else 0 for px in _pixels(image))
    mask = Image.frombytes("L", image.size, bits)
    ring = Image.frombytes(
        "L", image.size,
        bytes(255 if (d and not m) else 0
              for d, m in zip(_pixels(mask.filter(ImageFilter.MaxFilter(3))),
                              _pixels(mask))))
    counts: dict[tuple[int, int, int], int] = {}
    for px, on in zip(_pixels(image), _pixels(ring)):
        if on:
            counts[px] = counts.get(px, 0) + 1
    counts.pop(rgb, None)
    if not counts:
        return None
    return max(counts.items(), key=lambda kv: kv[1])[0]


def _dominant_colors(image, limit: int = 12) -> list[tuple[tuple[int, int, int], float]]:
    total = image.width * image.height
    counts = image.getcolors(maxcolors=total) or []
    counts.sort(key=lambda kv: kv[0], reverse=True)
    return [(rgb, n / total) for n, rgb in counts[:limit]]


def lint_render(path: Path, report: Report, label: str) -> None:
    try:
        from PIL import Image                   # type: ignore
    except ImportError:
        report.add("render.deps", "SKIP", label,
                   "Pillow not installed — no pixel metrics (pip install pillow)")
        return

    with Image.open(path) as probe:
        if probe.width < MIN_CANVAS_PX or probe.height < MIN_CANVAS_PX:
            report.add("render.scope", "SKIP", label,
                       f"{probe.width}×{probe.height} is below the canvas threshold — "
                       f"icons and logos are not judged as slides or pages")
            return

    with Image.open(path) as raw:
        image = raw.convert("RGB")
        # Quantise so anti-aliasing does not shatter the palette into thousands
        # of near-identical shades; keeps dominant-colour analysis meaningful.
        small = image.resize((min(image.width, 900), min(image.height, 900)))
        quant = small.quantize(colors=16, method=2).convert("RGB")

        total = image.width * image.height
        dark = 0
        for count, rgb in (image.getcolors(maxcolors=total) or []):
            if luminance(rgb) < 0.02:
                dark += count
        focal = dark / total

    if focal > FOCAL_COVERAGE_MAX:
        report.add("render.focal", "FAIL", label,
                   f"heavy surface covers {focal:.1%} of the canvas — it has stopped "
                   f"being emphasis and become the background",
                   measured=round(focal, 4), threshold=FOCAL_COVERAGE_MAX)
    else:
        report.add("render.focal", "PASS", label, f"heavy surface {focal:.1%}",
                   measured=round(focal, 4), threshold=FOCAL_COVERAGE_MAX)

    palette = _dominant_colors(quant)
    if not palette:
        report.add("render.contrast", "SKIP", label, "palette not readable")
        return
    background = palette[0][0]
    # Share alone cannot tell text from a tinted surface panel: a muted tile and a
    # paragraph can occupy the same fraction of the canvas, but only one of them
    # must clear a text-contrast floor.  Thinness separates them — glyph pixels sit
    # almost entirely on an edge, a solid panel is nearly all interior.
    detail = []
    for rgb, share in palette[1:]:
        if not 0.0005 <= share <= 0.15:
            continue
        thin = _thinness(quant, rgb)
        # thinness == 1.0 means the colour has no interior pixel at all: a 1px
        # scatter, which is what a renderer's anti-aliasing halo looks like.
        # Real glyph stems are 2px+ and always retain some interior.
        if TEXT_THINNESS_MIN <= thin < TEXT_THINNESS_MAX:
            detail.append((rgb, share, thin))
    if not detail:
        report.add("render.contrast", "SKIP", label,
                   "no text-like colour found (every minority colour reads as a solid "
                   "surface, not glyphs) — contrast unverified on pixels")
        return
    scored = []
    for rgb, share, thin in detail:
        backdrop = _local_backdrop(quant, rgb) or background
        scored.append((contrast_ratio(rgb, backdrop), rgb, share, thin, backdrop))
    worst, worst_rgb, worst_share, worst_thin, worst_bg = min(scored, key=lambda t: t[0])
    label = (f"dimmest text-like colour rgb{worst_rgb} ({worst_share:.2%} of canvas, "
             f"thinness {worst_thin:.2f}) on the surface behind it rgb{worst_bg} "
             f"= {worst:.2f}:1")
    # Deliberately capped at WARN.  Anti-aliasing fringe is thin and low-contrast by
    # construction, and no pixel-only test separates it from genuinely dim text.
    # A gate that cries wolf gets switched off, so the authoritative contrast lanes
    # are `html` (declared pairs) and the design system's own tokens; this probe
    # exists to point a human at the slide worth opening.
    if worst < CONTRAST_TEXT:
        report.add("render.contrast", "WARN", label,
                   f"{label} — under {CONTRAST_TEXT}:1. Heuristic: may be anti-aliasing "
                   f"fringe rather than text; confirm before treating as a defect",
                   measured=round(worst, 2), threshold=CONTRAST_TEXT)
    else:
        report.add("render.contrast", "PASS", label, label,
                   measured=round(worst, 2), threshold=CONTRAST_TEXT)


# ─── Lane: html (static analysis of artifact pages) ───


_STYLE_BLOCK = re.compile(r"<style[^>]*>(.*?)</style>", re.S | re.I)
_IMG_TAG = re.compile(r"<img\b[^>]*>", re.I)
_HEADING = re.compile(r"<h([1-6])\b", re.I)
_DECL = re.compile(r"([a-z-]+)\s*:\s*([^;{}]+)", re.I)
_RULE = re.compile(r"([^{}]+)\{([^{}]*)\}", re.S)
_FONT_SIZE_PX = re.compile(r"font-size\s*:\s*([\d.]+)px", re.I)


def lint_html(path: Path, report: Report, label: str) -> None:
    try:
        text = path.read_text(encoding="utf-8", errors="replace")
    except OSError as exc:
        report.add("html.read", "SKIP", label, f"unreadable: {exc}")
        return

    tag = label
    css = "\n".join(_STYLE_BLOCK.findall(text))
    inline_styles = re.findall(r'style\s*=\s*"([^"]*)"', text, re.I)
    all_css = css + "\n" + "\n".join(inline_styles)

    # -- motion respects preference ------------------------------------------
    has_motion = bool(re.search(r"(animation|transition)\s*:", all_css, re.I)) or \
        bool(re.search(r"@keyframes", css, re.I))
    guarded = bool(re.search(r"prefers-reduced-motion", text, re.I))
    if has_motion and not guarded:
        report.add("html.motion", "FAIL", tag,
                   "declares animation/transition with no prefers-reduced-motion guard")
    elif has_motion:
        report.add("html.motion", "PASS", tag, "motion is guarded by prefers-reduced-motion")

    # -- focus visibility -----------------------------------------------------
    kills_outline = re.search(r"outline\s*:\s*(none|0)\b", all_css, re.I)
    restores = re.search(r":focus(-visible)?\b[^{]*\{[^}]*(outline|box-shadow|border)", css, re.I)
    if kills_outline and not restores:
        report.add("html.focus", "FAIL", tag,
                   "removes the focus outline without providing a replacement indicator")
    elif kills_outline:
        report.add("html.focus", "PASS", tag, "outline removed but a focus style is restored")

    # -- images carry alt -----------------------------------------------------
    imgs = _IMG_TAG.findall(text)
    missing_alt = [t for t in imgs if not re.search(r"\balt\s*=", t, re.I)]
    if missing_alt:
        report.add("html.alt", "FAIL", tag,
                   f"{len(missing_alt)} of {len(imgs)} <img> tags have no alt attribute",
                   measured=len(missing_alt), threshold=0)
    elif imgs:
        report.add("html.alt", "PASS", tag, f"all {len(imgs)} <img> tags carry alt")

    # -- heading order --------------------------------------------------------
    levels = [int(m) for m in _HEADING.findall(text)]
    if levels:
        h1s = levels.count(1)
        problems = []
        if h1s == 0:
            problems.append("no <h1>")
        elif h1s > 1:
            problems.append(f"{h1s} <h1> elements")
        prev = None
        for lv in levels:
            if prev is not None and lv > prev + 1:
                problems.append(f"level jump h{prev}→h{lv}")
                break
            prev = lv
        if problems:
            report.add("html.headings", "WARN", tag,
                       "heading structure: " + ", ".join(problems))
        else:
            report.add("html.headings", "PASS", tag,
                       f"{len(levels)} headings, single h1, no skipped levels")

    # -- declared colour pairs ------------------------------------------------
    # A pair that cannot be statically resolved (unsupported colour function,
    # a custom property with no fallback) must never be silently dropped: if
    # every OTHER pair happens to pass, the unresolved one would otherwise
    # vanish into a bare PASS that never actually checked it.
    props = custom_properties(all_css)
    worst: tuple[float, str] | None = None
    unresolved: list[str] = []
    for selector, body in _RULE.findall(css):
        decls = {k.lower().strip(): v.strip() for k, v in _DECL.findall(body)}
        fg = decls.get("color")
        bg = decls.get("background-color") or decls.get("background")
        if not fg or not bg:
            continue
        sel = " ".join(selector.split())[:60]
        r_fg, r_bg = resolve_value(fg, props), resolve_value(bg, props)
        c_fg, c_bg = parse_color(r_fg), parse_color(r_bg)
        if not c_fg or not c_bg:
            unresolved.append(f"{sel} — {fg.strip()} on {bg.strip()}")
            continue
        ratio = contrast_ratio(c_fg, c_bg)
        if worst is None or ratio < worst[0]:
            worst = (ratio, f"{sel} — {r_fg.strip()} on {r_bg.strip()} = {ratio:.2f}:1")
    if worst is None:
        if unresolved:
            report.add("html.contrast", "SKIP", tag,
                       f"{len(unresolved)} rule(s) declare colour+background that could "
                       f"not be statically resolved (e.g. {unresolved[0]}) — unsupported "
                       f"colour function or a custom property with no fallback")
        else:
            report.add("html.contrast", "SKIP", tag,
                       "no rule declares colour and background together — static analysis "
                       "cannot resolve inherited or computed pairs")
    elif worst[0] < CONTRAST_TEXT:
        sev = "FAIL" if worst[0] < CONTRAST_LARGE else "WARN"
        report.add("html.contrast", sev, tag, worst[1] + " (AA text needs 4.5:1)",
                   measured=round(worst[0], 2), threshold=CONTRAST_TEXT)
    elif unresolved:
        report.add("html.contrast", "SKIP", tag,
                   f"worst resolvable pair clears {CONTRAST_TEXT}:1 ({worst[1]}), but "
                   f"{len(unresolved)} other pair(s) could not be statically resolved "
                   f"(e.g. {unresolved[0]}) and were never verified",
                   measured=round(worst[0], 2), threshold=CONTRAST_TEXT)
    else:
        report.add("html.contrast", "PASS", tag, "worst declared pair " + worst[1],
                   measured=round(worst[0], 2), threshold=CONTRAST_TEXT)

    # -- body type size -------------------------------------------------------
    sizes = [float(s) for s in _FONT_SIZE_PX.findall(all_css)]
    tiny = [s for s in sizes if s < BODY_MIN_PX]
    if tiny:
        report.add("html.type", "WARN", tag,
                   f"{len(tiny)} font-size declarations below {BODY_MIN_PX:.0f}px "
                   f"(smallest {min(tiny):.0f}px) — fine for metadata, not for body copy",
                   measured=min(tiny), threshold=BODY_MIN_PX)
    elif sizes:
        report.add("html.type", "PASS", tag,
                   f"{len(sizes)} font-size declarations, smallest {min(sizes):.0f}px")


# ─── Driver ───


IMAGE_SUFFIXES = {".png", ".jpg", ".jpeg", ".webp", ".bmp"}


def collect(target: Path, lane: str) -> list[Path]:
    if target.is_file():
        return [target]
    if lane == "render":
        return sorted(p for p in target.rglob("*") if p.suffix.lower() in IMAGE_SUFFIXES)
    if lane == "html":
        return sorted(target.rglob("*.html"))
    return sorted(target.rglob("*.pptx"))


def main(argv: list[str] | None = None) -> int:
    parser = argparse.ArgumentParser(
        description="Measure rendered design output against the Clade design floors.")
    parser.add_argument("lane", choices=["deck", "render", "html"])
    parser.add_argument("target", type=Path, help="file or directory")
    parser.add_argument("--json", action="store_true", help="emit findings as JSON")
    parser.add_argument("--quiet", action="store_true", help="only show FAIL/WARN")
    args = parser.parse_args(argv)

    if not args.target.exists():
        print(f"design-lint: no such path: {args.target}", file=sys.stderr)
        return 2

    paths = collect(args.target, args.lane)
    if not paths:
        print(f"design-lint: nothing to check under {args.target}", file=sys.stderr)
        return 2

    report = Report(lane=args.lane)
    runner = {"deck": lint_deck, "render": lint_render, "html": lint_html}[args.lane]
    root = args.target if args.target.is_dir() else args.target.parent
    for path in paths:
        # Label by path relative to the scan root. Artifact pages are all named
        # index.html, so a basename label silently merges hundreds of distinct
        # findings into one target and makes every aggregate count wrong.
        try:
            label = str(path.relative_to(root))
        except ValueError:
            label = path.name
        try:
            runner(path, report, label)
        except Exception as exc:  # a crashed check must be visible, never silent
            report.add(f"{args.lane}.error", "FAIL", label,
                       f"check crashed: {type(exc).__name__}: {exc}")

    if args.json:
        print(json.dumps({"lane": report.lane,
                          "failed": report.failed,
                          "findings": [asdict(f) for f in report.findings]}, indent=2))
    else:
        shown = report.findings
        if args.quiet:
            shown = [f for f in shown if f.severity in ("FAIL", "WARN")]
        filtered = Report(lane=report.lane, findings=shown)
        print(f"design-lint {args.lane}: {len(paths)} target(s)")
        print(filtered.render() if shown else "  (no findings)")
        if args.quiet and not shown:
            print(f"  ── {report.lane}: all checks clear")

    return 1 if report.failed else 0


if __name__ == "__main__":
    sys.exit(main())
